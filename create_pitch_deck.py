"""Generate a single-slide stakeholder pitch deck for the patient triage use case."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

NAVY = RGBColor(20, 45, 78)
BLUE = RGBColor(41, 98, 255)
TEAL = RGBColor(0, 137, 123)
AMBER = RGBColor(255, 179, 0)
RED = RGBColor(198, 40, 40)
LIGHT_BG = RGBColor(246, 248, 251)
CARD_BG = RGBColor(255, 255, 255)
CARD_BORDER = RGBColor(214, 220, 229)
TEXT = RGBColor(33, 37, 41)
MUTED = RGBColor(91, 104, 118)
WHITE = RGBColor(255, 255, 255)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=12,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT,
    vertical=MSO_VERTICAL_ANCHOR.TOP,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = vertical
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return textbox


def add_bullets(slide, left, top, width, height, items, font_size=10, color=TEXT):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(2)
    return textbox


def add_card(slide, left, top, width, height, title, title_color, bullets, title_size=13, bullet_size=10):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(1.2)

    add_textbox(
        slide,
        left + Inches(0.18),
        top + Inches(0.10),
        width - Inches(0.36),
        Inches(0.28),
        title,
        font_size=title_size,
        bold=True,
        color=title_color,
    )
    add_bullets(
        slide,
        left + Inches(0.18),
        top + Inches(0.42),
        width - Inches(0.32),
        height - Inches(0.52),
        bullets,
        font_size=bullet_size,
    )


def build_slide():
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        SLIDE_HEIGHT,
    )
    background.fill.solid()
    background.fill.fore_color.rgb = LIGHT_BG
    background.line.fill.background()

    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        SLIDE_WIDTH,
        Inches(0.95),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.fill.background()

    add_textbox(
        slide,
        Inches(0.35),
        Inches(0.14),
        Inches(8.4),
        Inches(0.36),
        "AI-Powered Patient Complaint and Incident Triage",
        font_size=24,
        bold=True,
        color=WHITE,
    )
    add_textbox(
        slide,
        Inches(0.35),
        Inches(0.50),
        Inches(8.7),
        Inches(0.22),
        "Stakeholder overview: faster intake, safer prioritization, clearer routing",
        font_size=11,
        color=RGBColor(212, 223, 236),
    )

    add_card(
        slide,
        Inches(0.35),
        Inches(1.2),
        Inches(4.0),
        Inches(1.9),
        "Problem Statement",
        RED,
        [
            "Complaint and incident review is manual, inconsistent, and slow.",
            "High-risk cases can sit in queues without timely escalation.",
            "Teams spend time sorting reports instead of acting on them.",
            "Leaders have limited visibility into trends, bottlenecks, and risk exposure.",
        ],
    )

    add_card(
        slide,
        Inches(4.65),
        Inches(1.2),
        Inches(4.0),
        Inches(1.9),
        "Benefits / Value",
        TEAL,
        [
            "Accelerates triage from hours or days to near real time.",
            "Improves consistency in classification and urgency decisions.",
            "Supports patient safety, service recovery, and compliance response.",
            "Reduces administrative burden and creates cleaner data for reporting.",
        ],
    )

    add_card(
        slide,
        Inches(8.95),
        Inches(1.2),
        Inches(4.0),
        Inches(1.9),
        "Solution Outputs",
        BLUE,
        [
            "Incident category and severity classification.",
            "Priority score and urgency flag for follow-up timing.",
            "Recommended routing to the right operational owner.",
            "Summary, root-cause cues, and immediate action guidance.",
        ],
    )

    flow_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.35),
        Inches(3.45),
        Inches(8.3),
        Inches(2.65),
    )
    flow_card.fill.solid()
    flow_card.fill.fore_color.rgb = CARD_BG
    flow_card.line.color.rgb = CARD_BORDER
    flow_card.line.width = Pt(1.2)

    add_textbox(
        slide,
        Inches(0.53),
        Inches(3.58),
        Inches(2.5),
        Inches(0.25),
        "How It Works",
        font_size=13,
        bold=True,
        color=AMBER,
    )

    steps = [
        ("1", "Capture", "Complaint or incident is submitted from any intake channel."),
        ("2", "Analyze", "LLM reviews narrative details, entities, and risk indicators."),
        ("3", "Prioritize", "System scores urgency based on harm, recurrence, and impact."),
        ("4", "Route", "Case is assigned or escalated to the right team with rationale."),
    ]
    x_positions = [0.55, 2.55, 4.55, 6.55]

    for index, (step_num, label, description) in enumerate(steps):
        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x_positions[index]),
            Inches(4.02),
            Inches(0.45),
            Inches(0.45),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE if index % 2 == 0 else TEAL
        circle.line.fill.background()

        add_textbox(
            slide,
            Inches(x_positions[index]),
            Inches(4.09),
            Inches(0.45),
            Inches(0.18),
            step_num,
            font_size=12,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            Inches(x_positions[index] + 0.58),
            Inches(3.95),
            Inches(1.15),
            Inches(0.22),
            label,
            font_size=12,
            bold=True,
            color=NAVY,
        )
        add_textbox(
            slide,
            Inches(x_positions[index] + 0.58),
            Inches(4.23),
            Inches(1.30),
            Inches(0.92),
            description,
            font_size=9,
            color=MUTED,
        )

        if index < len(steps) - 1:
            add_textbox(
                slide,
                Inches(x_positions[index] + 1.72),
                Inches(4.08),
                Inches(0.20),
                Inches(0.18),
                ">",
                font_size=18,
                bold=True,
                color=BLUE,
                align=PP_ALIGN.CENTER,
            )

    add_card(
        slide,
        Inches(8.95),
        Inches(3.45),
        Inches(4.0),
        Inches(2.65),
        "Next Steps",
        NAVY,
        [
            "Validate with a pilot set of historical complaints and incident reports.",
            "Define routing rules, review thresholds, and governance owners.",
            "Integrate with intake workflow and reporting dashboards.",
            "Measure impact on turnaround time, safety escalation, and staff effort.",
        ],
    )

    footer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(6.95),
        SLIDE_WIDTH,
        Inches(0.55),
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = NAVY
    footer.line.fill.background()

    add_textbox(
        slide,
        Inches(0.35),
        Inches(7.08),
        Inches(12.6),
        Inches(0.18),
        "Outcome: a scalable intake layer that helps teams identify the right case, at the right priority, with the right owner.",
        font_size=11,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    return presentation


def main():
    presentation = build_slide()
    output_path = Path(__file__).parent / "Patient_Triage_Stakeholder_One_Slide.pptx"
    presentation.save(output_path)
    print(f"Created {output_path}")


if __name__ == "__main__":
    main()
